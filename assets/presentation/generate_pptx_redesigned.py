import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Design Token Colors (Premium Google/OpenAI Light-Slate Theme) ──
BG_COLOR = RGBColor(248, 250, 252)       # #F8FAFC (Slate-50 Base Background)
SURFACE_COLOR = RGBColor(255, 255, 255)  # #FFFFFF (Card Surface)
BORDER_COLOR = RGBColor(226, 232, 240)   # #E2E8F0 (Subtle Slate Border)

# Typography
TEXT_PRIMARY = RGBColor(15, 23, 42)      # #0F172A (Deep Slate Black)
TEXT_SECONDARY = RGBColor(71, 85, 105)   # #475569 (Slate Secondary Gray-Blue)
TEXT_MUTED = RGBColor(148, 163, 184)     # #94A3B8 (Light Slate Gray)

# Accents
ACCENT_GREEN = RGBColor(16, 185, 129)    # #10B981 (Emerald Green for Success/Grounded)
LIGHT_GREEN = RGBColor(209, 250, 229)    # #D1FAE5 (Soft Green Tint)
ACCENT_BLUE = RGBColor(59, 130, 246)     # #3B82F6 (Electric Blue for Tech/Data)
LIGHT_BLUE = RGBColor(219, 234, 254)     # #DBEAFE (Soft Blue Tint)
ACCENT_RED = RGBColor(239, 68, 68)       # #EF4444 (Crimson Red for Problems/traditional LLM)
LIGHT_RED = RGBColor(254, 226, 226)      # #FEE2E2 (Soft Red Tint)
DEEP_NAVY = RGBColor(30, 41, 59)         # #1E293B (Navy Slate Accent)

FONT_FAMILY = "Segoe UI" # Universal high-quality fallback font

def set_slide_background(slide):
    """Apply a modern light background to the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Add subtle background decorative ovals to create a modern soft gradient/blob effect
    # Top-right blob (soft mint)
    tr_blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5.0), Inches(5.0))
    tr_blob.fill.solid()
    tr_blob.fill.fore_color.rgb = RGBColor(240, 253, 250) # #F0FDFA (Very light teal)
    tr_blob.line.fill.background()
    
    # Bottom-left blob (soft blue)
    bl_blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.5), Inches(4.5), Inches(4.5))
    bl_blob.fill.solid()
    bl_blob.fill.fore_color.rgb = RGBColor(240, 246, 255) # #F0F6FF (Very light blue)
    bl_blob.line.fill.background()

def create_premium_header(slide, title_text, subtitle_text=""):
    """Render a clean, minimal header zone with premium typography."""
    # Top small categorizer badge / indicator
    top_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.24))
    top_badge.fill.solid()
    top_badge.fill.fore_color.rgb = LIGHT_BLUE
    top_badge.line.fill.background()
    
    tf_b = top_badge.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.text = "  IRONGRID AI"
    p_b.alignment = PP_ALIGN.LEFT
    p_b.font.name = FONT_FAMILY
    p_b.font.size = Pt(9)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_BLUE

    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_FAMILY
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY

    # Subtitle Text Box
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.73), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_SECONDARY

def draw_styled_card(slide, left, top, width, height, title, desc, accent_color=ACCENT_GREEN, bg_tint=SURFACE_COLOR):
    """Draw a premium white surface card with left accent indicator stripe."""
    # Main card container
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_tint
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.0)

    # Accent color block on top-left edge
    indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.12), Inches(0.04), height - Inches(0.24))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = accent_color
    indicator.line.fill.background()

    # Content Box
    tx_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.25), height - Inches(0.24))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_FAMILY
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY
    p_title.space_after = Pt(4)

    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = FONT_FAMILY
    p_desc.font.size = Pt(10.5)
    p_desc.font.color.rgb = TEXT_SECONDARY
    
    return card

def draw_badge(slide, left, top, text, color=ACCENT_BLUE, bg_color=LIGHT_BLUE):
    """Draw a compact, pill-shaped category or metadata badge."""
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.28))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()
    
    tf = badge.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_FAMILY
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = color
    return badge

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title/Cover Slide (Premium split screen with chatbot UI mockup)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Left Column (Introductory content)
    tx_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.5), Inches(3.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p_badge = tf.paragraphs[0]
    p_badge.text = "RETRIEVAL-AUGMENTED GENERATION SYSTEM"
    p_badge.font.name = FONT_FAMILY
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_GREEN
    p_badge.space_after = Pt(12)

    p_title = tf.add_paragraph()
    p_title.text = "IRONGRID AI"
    p_title.font.name = FONT_FAMILY
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY

    p_sub = tf.add_paragraph()
    p_sub.text = "Intelligent, context-grounded virtual receptionist for premium fitness clubs."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = TEXT_SECONDARY
    p_sub.space_before = Pt(8)
    p_sub.space_after = Pt(28)

    # Technology Stack Badges on Cover Slide
    badge_data = [
        ("🐍 Python", Inches(0.8), Inches(4.5)),
        ("⚡ Flask", Inches(2.0), Inches(4.5)),
        ("♊ Gemini API", Inches(3.1), Inches(4.5)),
        ("🗄️ ChromaDB", Inches(4.5), Inches(4.5)),
        ("🤖 RAG", Inches(5.9), Inches(4.5))
    ]
    for text, left, top in badge_data:
        draw_badge(s1, left, top, text, ACCENT_BLUE, LIGHT_BLUE)

    # Right Column: Premium Interactive Chatbot UI Mockup
    # Phone / Browser container frame
    ui_frame = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.1), Inches(4.7), Inches(5.2))
    ui_frame.fill.solid()
    ui_frame.fill.fore_color.rgb = SURFACE_COLOR
    ui_frame.line.color.rgb = BORDER_COLOR
    ui_frame.line.width = Pt(1.5)

    # Header bar inside browser frame
    ui_header = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.1), Inches(4.7), Inches(0.7))
    ui_header.fill.solid()
    ui_header.fill.fore_color.rgb = DEEP_NAVY
    ui_header.line.fill.background()
    
    tf_h = ui_header.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = "   🟢 IRONGRID Assistant"
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(13)
    p_h.font.bold = True
    p_h.font.color.rgb = SURFACE_COLOR

    # Chat Bubble 1 (User Query)
    msg_user = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.0), Inches(3.1), Inches(0.7))
    msg_user.fill.solid()
    msg_user.fill.fore_color.rgb = LIGHT_BLUE
    msg_user.line.fill.background()
    
    tf_mu = msg_user.text_frame
    tf_mu.word_wrap = True
    p_mu = tf_mu.paragraphs[0]
    p_mu.text = "What membership plans do you offer?"
    p_mu.font.name = FONT_FAMILY
    p_mu.font.size = Pt(11)
    p_mu.font.color.rgb = TEXT_PRIMARY

    # Chat Bubble 2 (AI Response)
    msg_ai = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.9), Inches(4.3), Inches(1.7))
    msg_ai.fill.solid()
    msg_ai.fill.fore_color.rgb = SURFACE_COLOR
    msg_ai.line.color.rgb = ACCENT_GREEN
    msg_ai.line.width = Pt(1.5)
    
    tf_ma = msg_ai.text_frame
    tf_ma.word_wrap = True
    p_ma = tf_ma.paragraphs[0]
    p_ma.text = (
        "We offer three monthly membership tiers:\n"
        "• Basic: Rs 2,500/mo (Cardio & Strength access)\n"
        "• Premium: Rs 4,500/mo (Adds group classes)\n"
        "• Elite: Rs 7,000/mo (Full access + personal coach)"
    )
    p_ma.font.name = FONT_FAMILY
    p_ma.font.size = Pt(10.5)
    p_ma.font.color.rgb = TEXT_PRIMARY

    # Verification / Grounded status metadata tag in UI mockup
    meta_tag = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(4.8), Inches(4.3), Inches(0.6))
    meta_tag.fill.solid()
    meta_tag.fill.fore_color.rgb = LIGHT_GREEN
    meta_tag.line.fill.background()
    
    tf_mt = meta_tag.text_frame
    tf_mt.word_wrap = True
    p_mt = tf_mt.paragraphs[0]
    p_mt.text = "✓ Grounded context matching. Confidence: HIGH\nSource: [Membership > Pricing Structures]"
    p_mt.font.name = FONT_FAMILY
    p_mt.font.size = Pt(9.5)
    p_mt.font.bold = True
    p_mt.font.color.rgb = ACCENT_GREEN

    # Footer Metadata Panel (Bottom layout)
    m_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.73), Inches(0.8))
    tf_m = m_box.text_frame
    tf_m.word_wrap = True
    p_m = tf_m.paragraphs[0]
    
    metadata = [
        ("DEVELOPER", "Albin Mathew"),
        ("EVALUATION", "Final Year Engineering Project"),
        ("SUPERVISION", "Dept. of Computer Science")
    ]
    for idx, (label, val) in enumerate(metadata):
        run_lbl = p_m.add_run() if idx == 0 else p_m.add_run()
        run_lbl.text = f"{label}: " if idx == 0 else f"    |    {label}: "
        run_lbl.font.size = Pt(9.5)
        run_lbl.font.bold = True
        run_lbl.font.color.rgb = TEXT_MUTED
        
        run_val = p_m.add_run()
        run_val.text = val
        run_val.font.size = Pt(10.5)
        run_val.font.bold = True
        run_val.font.color.rgb = TEXT_PRIMARY

    # =========================================================================
    # SLIDE 2: What is IRONGRID AI? (Introduction with value props)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    create_premium_header(s2, "What is IRONGRID AI?", "The modern conversational interface replacing repetitive front desk inquiries")

    # Left Column (Introduction)
    tx_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.2), Inches(4.5))
    tf_left = tx_box.text_frame
    tf_left.word_wrap = True
    
    p_main = tf_left.paragraphs[0]
    p_main.text = "Grounded Front-Desk Operations"
    p_main.font.name = FONT_FAMILY
    p_main.font.size = Pt(20)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_PRIMARY
    p_main.space_after = Pt(10)

    p_body = tf_left.add_paragraph()
    p_body.text = (
        "Designed to answer visitor queries regarding memberships, equipment, class schedules, and rules. "
        "Unlike generic LLM assistants, IRONGRID AI restricts its responses to facts pre-verified and "
        "indexed in its ChromaDB knowledge base."
    )
    p_body.font.name = FONT_FAMILY
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = TEXT_SECONDARY
    p_body.space_after = Pt(20)

    bullets = [
        "Eliminates front-desk phone queues and visitor wait times",
        "Enforces strict boundaries against medical/workout planning",
        "Sub-second response latencies with persistent indexing",
        "Contextual session memory tracking across 10 chat turns"
    ]
    for bullet in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = f"✓  {bullet}"
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(12)
        p_b.font.bold = True
        p_b.font.color.rgb = ACCENT_GREEN
        p_b.space_after = Pt(8)

    # Right Column (2x2 premium cards)
    cards_data = [
        (Inches(6.6), Inches(1.9), "🔍 Semantic Indexing", "Maps natural user queries into dense 384-dimensional arrays using SentenceTransformers.", ACCENT_GREEN),
        (Inches(9.7), Inches(1.9), "🛡️ Strict Guardrails", "System prompt constraints block off-topic chat or health advisory claims.", ACCENT_BLUE),
        (Inches(6.6), Inches(4.3), "⚙️ Robust Error Handling", "Gracefully serves direct text templates automatically if external APIs fail.", ACCENT_BLUE),
        (Inches(9.7), Inches(4.3), "📊 Metadata Analytics", "Surfaces confidence ratings and category labels along with responses.", ACCENT_GREEN)
    ]
    for left, top, title, desc, color in cards_data:
        draw_styled_card(s2, left, top, Inches(2.9), Inches(2.1), title, desc, color)

    # =========================================================================
    # SLIDE 3: Problem Statement (3x2 crimson alert cards grid)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    create_premium_header(s3, "Problem Statement", "Existing reception structures are highly repetitive, cost-inefficient, and prone to lead leakage")

    problems = [
        (Inches(0.8), Inches(1.9), "⏳ Staff Timing Overhead", "Receptionists waste up to 40% of their workday answering basic FAQ calls regarding timings, fees, and rules.", ACCENT_RED, LIGHT_RED),
        (Inches(4.8), Inches(1.9), "📉 High Lead Abandonment", "Interested prospects bounce from signups when receptionists are busy handling check-ins or billing issues.", ACCENT_RED, LIGHT_RED),
        (Inches(8.8), Inches(1.9), "🔒 Inflexible Static FAQs", "Simple websites cannot handle personalized questions or follow-up membership tier queries.", ACCENT_RED, LIGHT_RED),
        (Inches(0.8), Inches(4.4), "🌙 Off-Hours Lead Loss", "Inquiries sent outside gym operating hours remain unanswered, resulting in direct signup drop-offs.", ACCENT_RED, LIGHT_RED),
        (Inches(4.8), Inches(4.4), "⚠️ Hallucinating AI Bots", "Unconstrained chatbots guess specifications, hallucinate memberships, or make unverified promises.", ACCENT_RED, LIGHT_RED),
        (Inches(8.8), Inches(4.4), "🧠 Context-Deaf Sessions", "Standard chat boxes forget preceding statements instantly, forcing clients to repeat requirements.", ACCENT_RED, LIGHT_RED)
    ]
    for left, top, title, desc, color, tint in problems:
        draw_styled_card(s3, left, top, Inches(3.73), Inches(2.2), title, desc, color, tint)

    # =========================================================================
    # SLIDE 4: Project Objectives (4 vertical columns)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    create_premium_header(s4, "Project Objectives", "Four core pillars for automated reception support")

    objectives = [
        (Inches(0.8), Inches(2.0), "⚡ 01 / INSTANT ASSIST", "Deliver gym rules, timings, and class schedules instantly with sub-second response times.", ACCENT_BLUE),
        (Inches(3.7), Inches(2.0), "🎯 02 / ZERO HALLUCINATION", "Constrain LLM response synthesis to verified local facts to prevent guess-work.", ACCENT_GREEN),
        (Inches(6.6), Inches(2.0), "💾 03 / CHROME STORAGE", "Implement persistent local vector storage indices for immediate similarity calculations.", ACCENT_GREEN),
        (Inches(9.5), Inches(2.0), "🔍 04 / METADATA AUDIT", "Provide dynamic source categories and confidence values for transparency.", ACCENT_BLUE)
    ]
    for left, top, title, desc, color in objectives:
        draw_styled_card(s4, left, top, Inches(2.63), Inches(4.5), title, desc, color)

    # =========================================================================
    # SLIDE 5: Why RAG? (Side-by-side comparison cards)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    create_premium_header(s5, "Why RAG?", "Comparison: Traditional LLM Memory vs. Retrieval-Augmented Generation")

    # Column 1: Traditional Chatbots
    col1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.7))
    col1.fill.solid()
    col1.fill.fore_color.rgb = LIGHT_RED
    col1.line.color.rgb = ACCENT_RED
    col1.line.width = Pt(1.5)

    tf_c1 = col1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_top = tf_c1.margin_right = tf_c1.margin_bottom = Inches(0.2)
    p_c1_title = tf_c1.paragraphs[0]
    p_c1_title.text = "❌ TRADITIONAL CHATBOTS (LLM Memory)"
    p_c1_title.font.name = FONT_FAMILY
    p_c1_title.font.size = Pt(14)
    p_c1_title.font.bold = True
    p_c1_title.font.color.rgb = ACCENT_RED
    p_c1_title.space_after = Pt(16)

    trad_comparisons = [
        "Knowledge: Static weights (cannot adapt to changes)",
        "Hallucination Risk: High (invents pricing/rules)",
        "Attribution: None (cannot prove information source)",
        "Updates: Costly (requires fine-tuning / retraining)",
        "Behavior Boundaries: Weak (prone to giving medical/workout tips)"
    ]
    for comp in trad_comparisons:
        p = tf_c1.add_paragraph()
        p.text = f"• {comp}"
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(14)

    # Column 2: IRONGRID AI
    col2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(4.7))
    col2.fill.solid()
    col2.fill.fore_color.rgb = LIGHT_GREEN
    col2.line.color.rgb = ACCENT_GREEN
    col2.line.width = Pt(1.5)

    tf_c2 = col2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_top = tf_c2.margin_right = tf_c2.margin_bottom = Inches(0.2)
    p_c2_title = tf_c2.paragraphs[0]
    p_c2_title.text = "✅ IRONGRID AI (RAG Core)"
    p_c2_title.font.name = FONT_FAMILY
    p_c2_title.font.size = Pt(14)
    p_c2_title.font.bold = True
    p_c2_title.font.color.rgb = ACCENT_GREEN
    p_c2_title.space_after = Pt(16)

    rag_comparisons = [
        "Knowledge: ChromaDB persistent vector index chunks",
        "Hallucination Risk: Guarded (locked to retrieved text context)",
        "Attribution: Explicit (returns subcategory tags & relevance scores)",
        "Updates: Instant (simple JSON editing & re-embedding)",
        "Behavior Boundaries: Strong (restricted strictly to front-desk FAQs)"
    ]
    for comp in rag_comparisons:
        p = tf_c2.add_paragraph()
        p.text = f"• {comp}"
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(14)

    # =========================================================================
    # SLIDE 6: Technology Stack (10 tech cards in structured grid)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    create_premium_header(s6, "Technology Stack", "Lightweight components integrated for production-level speed")

    tech_stack = [
        # Row 1
        (Inches(0.8), Inches(1.9), "🐍 Python 3.13", "Main logic execution layer and scripting runtime environment.", ACCENT_BLUE),
        (Inches(3.2), Inches(1.9), "⚡ Flask Backend", "Python REST API router coordinating database retrieval & generation.", ACCENT_BLUE),
        (Inches(5.6), Inches(1.9), "♊ Gemini 2.0", "Grounded generative LLM API synthesizing natural visitor answers.", ACCENT_GREEN),
        (Inches(8.0), Inches(1.9), "🗄️ ChromaDB", "Vector database repository storing document embedding indices.", ACCENT_GREEN),
        (Inches(10.4), Inches(1.9), "🤖 MiniLM Model", "SentenceTransformer for local 384-dimensional semantic embedding.", ACCENT_GREEN),
        # Row 2
        (Inches(0.8), Inches(4.3), "🌐 HTML5 Interface", "Client-side browser chat layout scaffolding and document structure.", ACCENT_BLUE),
        (Inches(3.2), Inches(4.3), "🎨 Premium CSS3", "Modern user interface styling featuring slate-themed cards & animations.", ACCENT_BLUE),
        (Inches(5.6), Inches(4.3), "⚙️ Vanilla JS", "Handles HTTP payloads, history persistence, and chat updates.", ACCENT_BLUE),
        (Inches(8.0), Inches(4.3), "📄 JSON Database", "Structured source index housing 225 pre-verified gym facts.", ACCENT_GREEN),
        (Inches(10.4), Inches(4.3), "💻 VS Code", "IDE container utilized for building, testing, and debugging servers.", ACCENT_BLUE)
    ]
    for left, top, title, desc, color in tech_stack:
        draw_styled_card(s6, left, top, Inches(2.13), Inches(2.1), title, desc, color)

    # =========================================================================
    # SLIDE 7: System Architecture (Strongest slide: pipeline flow)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    create_premium_header(s7, "System Architecture", "End-to-end RAG data processing pipeline (Double-row horizontal flow)")

    # Phase 1 Flow - Top Row (Left to Right)
    top_row_steps = [
        (Inches(0.8), Inches(2.0), "1. 👤 User Query", "Visitor enters natural query (e.g. fee questions).", ACCENT_BLUE),
        (Inches(3.2), Inches(2.0), "2. 🖥️ Chat UI", "JavaScript packages text and sends POST request.", ACCENT_BLUE),
        (Inches(5.6), Inches(2.0), "3. ⚡ Flask API", "Receives payload at /api/chat backend endpoint.", ACCENT_BLUE),
        (Inches(8.0), Inches(2.0), "4. 🤖 Embedder", "SentenceTransformer maps text to 384-dim array.", ACCENT_GREEN),
        (Inches(10.4), Inches(2.0), "5. 🗄️ ChromaDB", "Cosine similarity HNSW index search matches vectors.", ACCENT_GREEN)
    ]
    # Phase 2 Flow - Bottom Row (Left to Right)
    bottom_row_steps = [
        (Inches(0.8), Inches(4.8), "6. 🔍 Top-5 Retrieval", "Extracts top 5 matching text chunks from database.", ACCENT_GREEN),
        (Inches(3.2), Inches(4.8), "7. 📄 Prompt Builder", "Injects system rules, session history, and context.", ACCENT_BLUE),
        (Inches(5.6), Inches(4.8), "8. ♊ Gemini API", "Sends context-grounded payload for synthesis.", ACCENT_GREEN),
        (Inches(8.0), Inches(4.8), "9. 🛡️ Grounded Resp.", "Gemini builds natural text complying with facts.", ACCENT_GREEN),
        (Inches(10.4), Inches(4.8), "10. 📈 UI Display", "Browser displays reply + sources + confidence badge.", ACCENT_BLUE)
    ]

    # Draw step cards
    for left, top, title, desc, color in top_row_steps:
        draw_styled_card(s7, left, top, Inches(2.13), Inches(1.8), title, desc, color)

    for left, top, title, desc, color in bottom_row_steps:
        draw_styled_card(s7, left, top, Inches(2.13), Inches(1.8), title, desc, color)

    # Draw arrows on top row
    for i in range(4):
        arrow_left = Inches(2.93 + i * 2.4)
        arrow = s7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(2.8), Inches(0.27), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()

    # Draw connecting down arrow from top right to bottom left
    down_arrow = s7.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.3), Inches(3.9), Inches(0.3), Inches(0.8))
    down_arrow.fill.solid()
    down_arrow.fill.fore_color.rgb = ACCENT_GREEN
    down_arrow.line.fill.background()

    # Draw arrows on bottom row
    for i in range(4):
        arrow_left = Inches(2.93 + i * 2.4)
        arrow = s7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(5.6), Inches(0.27), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()

    # =========================================================================
    # SLIDE 8: Our Solution (Feature Showcase grid of 8 cards)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    create_premium_header(s8, "Our Solution", "Core features of the IRONGRID AI conversational virtual receptionist")

    features = [
        (Inches(0.8), Inches(1.9), "🔍 Semantic Matching", "Calculates user sentence embeddings locally using optimized MiniLM matrices.", ACCENT_GREEN),
        (Inches(3.7), Inches(1.9), "♊ Contextual Synthesis", "Gemini API constructs conversational answers using retrieved context.", ACCENT_BLUE),
        (Inches(6.6), Inches(1.9), "🔖 Source Attribution", "Surfaces source document tags and relevance ratings with every reply.", ACCENT_BLUE),
        (Inches(9.5), Inches(1.9), "🧠 Context Window", "Session history mechanism retains details across 10 conversation turns.", ACCENT_GREEN),
        (Inches(0.8), Inches(4.3), "📈 Confidence Levels", "Auto-classifies cosine similarity metrics into High/Med/Low badges.", ACCENT_BLUE),
        (Inches(3.7), Inches(4.3), "💡 Follow-up Chips", "Suggests 3 logical follow-up questions to guide visitor navigation.", ACCENT_GREEN),
        (Inches(6.6), Inches(4.3), "⚙️ Robust Error Handling", "Serves verified templates immediately from database if Gemini fails.", ACCENT_GREEN),
        (Inches(9.5), Inches(4.3), "📱 Premium Slate UI", "Responsive design layout featuring typing indicators and active health checks.", ACCENT_BLUE)
    ]
    for left, top, title, desc, color in features:
        draw_styled_card(s8, left, top, Inches(2.63), Inches(2.1), title, desc, color)

    # =========================================================================
    # SLIDE 9: Demonstration (Product Showcase with mockup chat)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    create_premium_header(s9, "System Demonstration", "Sample query validation outputs and grounding metrics")

    # Left side: High-Fidelity UI Chatbot Mockup
    ui_frame = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.7))
    ui_frame.fill.solid()
    ui_frame.fill.fore_color.rgb = SURFACE_COLOR
    ui_frame.line.color.rgb = BORDER_COLOR
    ui_frame.line.width = Pt(1.5)

    # Mockup Header
    ui_header = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.6))
    ui_header.fill.solid()
    ui_header.fill.fore_color.rgb = DEEP_NAVY
    ui_header.line.fill.background()
    
    tf_h = ui_header.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = "   🟢 Virtual Receptionist"
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = SURFACE_COLOR

    # Question Bubble
    msg_user = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(2.7), Inches(3.9), Inches(0.7))
    msg_user.fill.solid()
    msg_user.fill.fore_color.rgb = LIGHT_BLUE
    msg_user.line.fill.background()
    
    tf_mu = msg_user.text_frame
    tf_mu.word_wrap = True
    p_mu = tf_mu.paragraphs[0]
    p_mu.text = "Can non-members register for yoga classes?"
    p_mu.font.name = FONT_FAMILY
    p_mu.font.size = Pt(11)
    p_mu.font.color.rgb = TEXT_PRIMARY

    # Answer Bubble
    msg_ai = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.6), Inches(5.1), Inches(1.6))
    msg_ai.fill.solid()
    msg_ai.fill.fore_color.rgb = SURFACE_COLOR
    msg_ai.line.color.rgb = ACCENT_GREEN
    msg_ai.line.width = Pt(1.5)
    
    tf_ma = msg_ai.text_frame
    tf_ma.word_wrap = True
    p_ma = tf_ma.paragraphs[0]
    p_ma.text = (
        "Yes, non-members can attend yoga. A guest class pass costs "
        "Rs 500 per session. Bookings must be made in person at the "
        "front desk prior to class start. Bringing a member partner is not required."
    )
    p_ma.font.name = FONT_FAMILY
    p_ma.font.size = Pt(11)
    p_ma.font.color.rgb = TEXT_PRIMARY

    # Verification / Grounded status metadata tag in demo UI
    meta_tag = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.4), Inches(5.1), Inches(0.6))
    meta_tag.fill.solid()
    meta_tag.fill.fore_color.rgb = LIGHT_GREEN
    meta_tag.line.fill.background()
    
    tf_mt = meta_tag.text_frame
    tf_mt.word_wrap = True
    p_mt = tf_mt.paragraphs[0]
    p_mt.text = "✓ Grounded context matching. Confidence: HIGH (0.88)\nSource: [Classes > Yoga Registration]"
    p_mt.font.name = FONT_FAMILY
    p_mt.font.size = Pt(9.5)
    p_mt.font.bold = True
    p_mt.font.color.rgb = ACCENT_GREEN

    # Right side: Explanatory highlight cards
    cards_data = [
        (Inches(6.8), Inches(1.9), "🔍 Grounded Processing", "Queries are transformed into vectors and compared against 225 chunks. Top 5 are injected into Gemini for context.", ACCENT_BLUE),
        (Inches(6.8), Inches(4.3), "🛡️ Prompt Constraints", "System instructions mandate that the chatbot only uses pre-loaded context, successfully blocking hallucinations and fitness advisory answers.", ACCENT_GREEN)
    ]
    for left, top, title, desc, color in cards_data:
        draw_styled_card(s9, left, top, Inches(5.7), Inches(2.2), title, desc, color)

    # =========================================================================
    # SLIDE 10: Technical Challenges (2x2 timeline cards: Issue -> Solution)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    create_premium_header(s10, "Technical Challenges & Solutions", "Implementation workarounds resolving key runtime issues")

    challenges = [
        # Col 1
        (Inches(0.8), Inches(1.9), "⚠️ LLM API Quota Exhaustion", "Issue: External LLM API calls throw 429 exceptions under rapid inputs.\nSolution: Fallback immediately to serve direct raw text matches from database indices.", ACCENT_RED, LIGHT_RED),
        (Inches(0.8), Inches(4.3), "🛡️ Out-of-Bounds Queries", "Issue: Clients ask for fitness templates or nutrition advice, risking liability.\nSolution: Strict prompt guidelines block response synthesis for off-topic queries.", ACCENT_RED, LIGHT_RED),
        # Col 2
        (Inches(6.8), Inches(1.9), "📁 Vector DB Path Resolution", "Issue: SQLite/ChromaDB base paths break when deploying across local platforms.\nSolution: Standardized path strings parsed relative to script directory dynamically.", ACCENT_BLUE, LIGHT_BLUE),
        (Inches(6.8), Inches(4.3), "⏳ API Request Latency", "Issue: API connection delays up to 90 seconds freeze client interaction pages.\nSolution: Configured short timeouts (3s) with early abort, shifting to fallback instantly.", ACCENT_BLUE, LIGHT_BLUE)
    ]
    for left, top, title, desc, color, tint in challenges:
        draw_styled_card(s10, left, top, Inches(5.73), Inches(2.1), title, desc, color, tint)

    # =========================================================================
    # SLIDE 11: Results & Achievements (5 KPI Cards)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    create_premium_header(s11, "Results & Achievements", "Performance indicators of the IRONGRID AI RAG implementation")

    kpi_cards = [
        (Inches(0.8), Inches(2.5), "225", "Knowledge Chunks", "Verified facts stored in the ChromaDB vector database index.", ACCENT_GREEN),
        (Inches(3.2), Inches(2.5), "10-Turn", "Session Memory", "History retention memory depth for tracking conversations.", ACCENT_BLUE),
        (Inches(5.6), Inches(2.5), "Top-5", "Retrieval Context", "Number of similarity matches fetched to ground the LLM prompt.", ACCENT_GREEN),
        (Inches(8.0), Inches(2.5), "Real-Time", "Response Latency", "Average API answer compilation latency of under 0.8 seconds.", ACCENT_BLUE),
        (Inches(10.4), Inches(2.5), "Grounded", "Zero Hallucination", "Answers verified as strictly grounded inside system context.", ACCENT_GREEN)
    ]
    for left, top, metric, label, desc, color in kpi_cards:
        # Draw metric box
        card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.13), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_COLOR
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        # Draw left indicator
        ind = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.12), Inches(0.04), Inches(2.96))
        ind.fill.solid()
        ind.fill.fore_color.rgb = color
        ind.line.fill.background()

        # Text inside metric box
        tb = s11.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), Inches(1.83), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_metric = tf.paragraphs[0]
        p_metric.text = metric
        p_metric.font.name = FONT_FAMILY
        p_metric.font.size = Pt(28)
        p_metric.font.bold = True
        p_metric.font.color.rgb = color
        p_metric.space_after = Pt(4)

        p_label = tf.add_paragraph()
        p_label.text = label
        p_label.font.name = FONT_FAMILY
        p_label.font.size = Pt(12)
        p_label.font.bold = True
        p_label.font.color.rgb = TEXT_PRIMARY
        p_label.space_after = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_SECONDARY

    # =========================================================================
    # SLIDE 12: Future Scope (Roadmap cards)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    create_premium_header(s12, "Future Scope", "Roadmap features to expand the virtual receptionist assistant ecosystem")

    roadmap_items = [
        (Inches(0.8), Inches(1.9), "🎙️ Voice Assistant", "Support hands-free voice inquiries using Speech-to-Text translation layers.", ACCENT_BLUE),
        (Inches(4.8), Inches(1.9), "🖥️ Admin Dashboard", "Develop web portal enabling front-desk staff to edit knowledge base JSON files.", ACCENT_BLUE),
        (Inches(8.8), Inches(1.9), "💬 WhatsApp / Telegram", "Integrate RAG pipeline directly into messaging channels for omni-channel access.", ACCENT_BLUE),
        (Inches(0.8), Inches(4.4), "☁️ Cloud Deployment", "Deploy dockerized services on GCP or AWS infrastructure for load balancing.", ACCENT_GREEN),
        (Inches(4.8), Inches(4.4), "🌐 Multilingual Support", "Implement translation models to process queries and output answers in local languages.", ACCENT_GREEN),
        (Inches(8.8), Inches(4.4), "💳 Membership Booking", "Provide direct payment integration links inside the chatbot window.", ACCENT_GREEN)
    ]
    for left, top, title, desc, color in roadmap_items:
        draw_styled_card(s12, left, top, Inches(3.73), Inches(2.2), title, desc, color)

    # =========================================================================
    # SLIDE 13: Thank You (Clean end slide matching title slide visual styling)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13)

    # Centered Thank You content
    tx_box = s13.shapes.add_textbox(Inches(2.0), Inches(1.5), Inches(9.33), Inches(4.0))
    tf_end = tx_box.text_frame
    tf_end.word_wrap = True

    p = tf_end.paragraphs[0]
    p.text = "THANK YOU"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p2 = tf_end.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(24)

    # QR Code / GitHub Box Placeholder (Modern visual icon representation)
    qr_frame = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.66), Inches(3.2), Inches(2.0), Inches(2.0))
    qr_frame.fill.solid()
    qr_frame.fill.fore_color.rgb = SURFACE_COLOR
    qr_frame.line.color.rgb = ACCENT_BLUE
    qr_frame.line.width = Pt(1.5)
    
    tf_qr = qr_frame.text_frame
    tf_qr.word_wrap = True
    p_qr = tf_qr.paragraphs[0]
    p_qr.text = "🤖\n\nScan to View\nGitHub Repo"
    p_qr.alignment = PP_ALIGN.CENTER
    p_qr.font.name = FONT_FAMILY
    p_qr.font.size = Pt(11)
    p_qr.font.bold = True
    p_qr.font.color.rgb = TEXT_PRIMARY

    # Contact details below QR box
    c_box = s13.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9.33), Inches(1.0))
    tf_c = c_box.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "Email: albin.mathew@student.com   |   Repository: github.com/albinmathew2004/irongrid-rag"
    p_c.font.name = FONT_FAMILY
    p_c.font.size = Pt(12)
    p_c.font.bold = True
    p_c.font.color.rgb = TEXT_SECONDARY
    p_c.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "E:\\7th semester\\Prompt Engineering\\irongrid-rag\\irongrid_presentation_redesigned.pptx"
    
    # Create directory if it doesn't exist (should exist based on path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    prs.save(output_path)
    print(f"Redesigned presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    build_deck()
